import os
import uuid
import logging
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import PyPDF2
from io import BytesIO
import json
import re
import io
from urllib.parse import urlparse, urljoin
import asyncio

# FastAPI and Pydantic imports
from pydantic import BaseModel

# LangChain imports
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.embeddings import HuggingFaceEmbeddings
    from langchain.schema import Document
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain.memory import ConversationBufferMemory
    from langchain.chains import ConversationalRetrievalChain
    from pinecone import Pinecone, ServerlessSpec
    LANGCHAIN_AVAILABLE = True
except ImportError as e:
    print(f"⚠️  LangChain components not available: {e}")
    LANGCHAIN_AVAILABLE = False

# Supabase imports
try:
    from supabase import create_client, Client
    SUPABASE_AVAILABLE = True
except ImportError as e:
    print(f"⚠️  Supabase client not available: {e}")
    SUPABASE_AVAILABLE = False

# Web scraping imports
try:
    import requests
    from firecrawl import FirecrawlApp
    WEB_SCRAPING_AVAILABLE = True
except ImportError as e:
    print(f"⚠️  Web scraping libraries not available: {e}")
    WEB_SCRAPING_AVAILABLE = False

# Document processing imports
try:
    from docx import Document as DocxDocument
    import openpyxl
    import xlrd
    from pptx import Presentation
    import markdown
    import csv
    from PIL import Image
    import pytesseract
    import speech_recognition as sr
    from pydub import AudioSegment
    import magic
    import chardet
    import validators
    DOCUMENT_PROCESSING_AVAILABLE = True
except ImportError as e:
    print(f"⚠️  Document processing libraries not available: {e}")
    DOCUMENT_PROCESSING_AVAILABLE = False

from models import ChatMessage, ChatResponse, Book

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class RAGService:
    def __init__(self):
        self.supabase_client = None
        self.pinecone_client = None
        self.pinecone_index = None
        self.embeddings = None
        self.llm = None
        self.text_splitter = None
        self.initialized = False
    
    async def initialize(self):
        """Initialize RAG components"""
        try:
            logger.info("Initializing RAG service...")
            
            # Initialize Supabase client
            if SUPABASE_AVAILABLE:
                supabase_url = os.getenv("SUPABASE_URL")
                supabase_key = os.getenv("SUPABASE_KEY")
                
                if supabase_url and supabase_key:
                    self.supabase_client = create_client(supabase_url, supabase_key)
                    logger.info("✅ Supabase client initialized")
                else:
                    logger.warning("⚠️  SUPABASE_URL or SUPABASE_KEY not found")
            
            if not LANGCHAIN_AVAILABLE:
                logger.warning("LangChain not available, using simple responses")
                self.initialized = True
                return
            
            # Initialize embeddings
            self.embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-mpnet-base-v2",
                model_kwargs={'device': 'cpu'}
            )
            
            # Initialize text splitter
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len
            )
            
            # Initialize Gemini LLM
            gemini_key = os.getenv("GEMINI_API_KEY")
            if gemini_key:
                self.llm = ChatGoogleGenerativeAI(
                    model="gemini-2.0-flash",
                    google_api_key=gemini_key,
                    temperature=0.3
                )
                logger.info("✅ Gemini 2.0 Flash LLM initialized")
            
            # Initialize Pinecone
            pinecone_key = os.getenv("PINECONE_API_KEY")
            pinecone_host = os.getenv("PINECONE_HOST")
            
            if pinecone_key:
                self.pinecone_client = Pinecone(api_key=pinecone_key)
                
                index_name = os.getenv("PINECONE_INDEX_NAME", "tutor-rag-index")
                
                # Check if index exists
                existing_indexes = [index.name for index in self.pinecone_client.list_indexes()]
                
                if index_name not in existing_indexes:
                    logger.info(f"Creating Pinecone index: {index_name}")
                    self.pinecone_client.create_index(
                        name=index_name,
                        dimension=768,  # all-mpnet-base-v2 dimension
                        metric="cosine",
                        spec=ServerlessSpec(
                            cloud="aws",
                            region="us-east-1"
                        )
                    )
                
                # Connect to index
                if pinecone_host:
                    self.pinecone_index = self.pinecone_client.Index(index_name, host=pinecone_host)
                    logger.info(f"✅ Pinecone connected with host: {pinecone_host}")
                else:
                    self.pinecone_index = self.pinecone_client.Index(index_name)
                    logger.info(f"✅ Pinecone connected")
            
            self.initialized = True
            logger.info("🚀 RAG service initialized successfully")
            
        except Exception as e:
            logger.error(f"❌ Error initializing RAG service: {e}")
            self.initialized = True  # Continue with limited functionality

    def get_file_type(self, filename: str, content: bytes = None) -> str:
        """Determine file type from filename and content"""
        extension = Path(filename).suffix.lower()
        
        SUPPORTED_FORMATS = {
            'pdf': ['.pdf'],
            'word': ['.docx', '.doc'],
            'excel': ['.xlsx', '.xls', '.csv'],
            'powerpoint': ['.pptx', '.ppt'],
            'text': ['.txt', '.md', '.markdown', '.json'],
            'image': ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'],
            'audio': ['.wav', '.mp3', '.m4a', '.flac', '.aac'],
            'web': ['.html', '.htm']
        }
        
        for format_type, extensions in SUPPORTED_FORMATS.items():
            if extension in extensions:
                return format_type
        
        # Try to detect MIME type if extension is unknown
        if content:
            try:
                import magic
                mime_type = magic.from_buffer(content, mime=True)
                if mime_type.startswith('text/'):
                    return 'text'
                elif mime_type.startswith('image/'):
                    return 'image'
                elif mime_type.startswith('audio/'):
                    return 'audio'
                elif 'pdf' in mime_type:
                    return 'pdf'
            except Exception:
                pass
        
        return 'text'  # Default fallback
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported"""
        extension = Path(filename).suffix.lower()
        supported_extensions = [
            '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.csv', 
            '.pptx', '.ppt', '.txt', '.md', '.markdown', '.json',
            '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif',
            '.wav', '.mp3', '.m4a', '.flac', '.aac', '.html', '.htm'
        ]
        return extension in supported_extensions

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file"""
        try:
            pdf_file = BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            if len(pdf_reader.pages) == 0:
                raise ValueError("PDF file contains no pages")
            
            text = ""
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num + 1} ---\n"
                        text += page_text
                except Exception as e:
                    logger.warning(f"Could not extract text from page {page_num + 1}: {e}")
                    continue
            
            text = text.strip()
            if not text:
                raise ValueError("No text could be extracted from the PDF")
            
            logger.info(f"Successfully extracted {len(text)} characters from {len(pdf_reader.pages)} pages")
            return text
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            raise

    def extract_text_from_document(self, file_content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from various document formats"""
        file_type = self.get_file_type(filename, file_content)
        metadata = {
            'filename': filename,
            'file_type': file_type,
            'file_size': len(file_content)
        }
        
        try:
            if file_type == 'pdf':
                text = self.extract_text_from_pdf(file_content)
                pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
                metadata.update({
                    'page_count': len(pdf_reader.pages),
                    'format_details': 'PDF Document'
                })
            elif file_type == 'word' and DOCUMENT_PROCESSING_AVAILABLE:
                text, word_metadata = self._extract_word_text(file_content)
                metadata.update(word_metadata)
            elif file_type == 'excel' and DOCUMENT_PROCESSING_AVAILABLE:
                text, excel_metadata = self._extract_excel_text(file_content, filename)
                metadata.update(excel_metadata)
            elif file_type == 'powerpoint' and DOCUMENT_PROCESSING_AVAILABLE:
                text, ppt_metadata = self._extract_powerpoint_text(file_content)
                metadata.update(ppt_metadata)
            elif file_type == 'text':
                text, text_metadata = self._extract_text_content(file_content, filename)
                metadata.update(text_metadata)
            elif file_type == 'image' and DOCUMENT_PROCESSING_AVAILABLE:
                text, img_metadata = self._extract_image_text(file_content)
                metadata.update(img_metadata)
            elif file_type == 'web':
                text, web_metadata = self._extract_html_text(file_content)
                metadata.update(web_metadata)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {e}")
            raise

    def _extract_word_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Word documents"""
        try:
            doc_file = BytesIO(content)
            doc = DocxDocument(doc_file)
            
            text = ""
            paragraph_count = 0
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
                    paragraph_count += 1
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text += cell.text + "\n"
            
            metadata = {
                'paragraph_count': paragraph_count,
                'table_count': len(doc.tables),
                'format_details': 'Word Document'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing Word document: {e}")

    def _extract_excel_text(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Excel files"""
        try:
            extension = Path(filename).suffix.lower()
            
            if extension == '.csv':
                return self._extract_csv_text(content)
            
            text = ""
            sheet_count = 0
            
            if extension == '.xlsx':
                workbook = openpyxl.load_workbook(BytesIO(content), data_only=True)
                sheet_names = workbook.sheetnames
                
                for sheet_name in sheet_names:
                    sheet = workbook[sheet_name]
                    text += f"\n--- Sheet: {sheet_name} ---\n"
                    
                    for row in sheet.iter_rows():
                        row_text = []
                        for cell in row:
                            if cell.value is not None:
                                row_text.append(str(cell.value))
                        if row_text:
                            text += " | ".join(row_text) + "\n"
                    
                    sheet_count += 1
            
            metadata = {
                'sheet_count': sheet_count,
                'format_details': f'Excel Spreadsheet ({extension})'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing Excel file: {e}")

    def _extract_csv_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from CSV files"""
        try:
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            
            text_content = content.decode(encoding)
            csv_reader = csv.reader(io.StringIO(text_content))
            
            text = ""
            row_count = 0
            
            for row in csv_reader:
                if row:
                    text += " | ".join(str(cell) for cell in row) + "\n"
                    row_count += 1
            
            metadata = {
                'row_count': row_count,
                'encoding': encoding,
                'format_details': 'CSV File'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing CSV file: {e}")

    def _extract_powerpoint_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PowerPoint presentations"""
        try:
            ppt_file = BytesIO(content)
            prs = Presentation(ppt_file)
            
            text = ""
            slide_count = 0
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                
                slide_count += 1
            
            metadata = {
                'slide_count': slide_count,
                'format_details': 'PowerPoint Presentation'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing PowerPoint file: {e}")

    def _extract_text_content(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from text-based files"""
        try:
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            
            text = content.decode(encoding)
            extension = Path(filename).suffix.lower()
            
            metadata = {
                'encoding': encoding,
                'character_count': len(text),
                'line_count': len(text.splitlines()),
                'format_details': f'Text File ({extension})'
            }
            
            if extension in ['.md', '.markdown']:
                metadata['format_details'] = 'Markdown Document'
            elif extension == '.json':
                try:
                    json_data = json.loads(text)
                    text = json.dumps(json_data, indent=2)
                    metadata['format_details'] = 'JSON Document'
                except json.JSONDecodeError:
                    pass
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing text file: {e}")

    def _extract_image_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from images using OCR"""
        try:
            image = Image.open(BytesIO(content))
            
            metadata = {
                'image_size': image.size,
                'image_mode': image.mode,
                'format_details': f'Image with OCR ({image.format})'
            }
            
            text = pytesseract.image_to_string(image)
            
            if not text.strip():
                text = "[No text found in image]"
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing image: {e}")

    def _extract_html_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from HTML files"""
        try:
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            
            html_content = content.decode(encoding)
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Convert to markdown-like text
            h = html2text.HTML2Text()
            h.ignore_links = False
            text = h.handle(html_content)
            
            metadata = {
                'encoding': encoding,
                'title': soup.title.string if soup.title else 'No title',
                'format_details': 'HTML Document'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing HTML file: {e}")

    async def process_book(self, file_content: bytes, filename: str, user_id: str) -> str:
        """Process a document and store it (supports multiple formats)"""
        try:
            book_id = str(uuid.uuid4())
            logger.info(f"Processing document: {filename} for user: {user_id}")
            
            # Check if format is supported
            if not self.is_supported_format(filename):
                raise ValueError(f"Unsupported file format: {filename}")
            
            # Extract text from document (multiple formats supported)
            text, metadata = self.extract_text_from_document(file_content, filename)
            
            # Store book in Supabase
            if self.supabase_client:
                try:
                    # Prepare book data with metadata
                    title = filename
                    for ext in ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt', '.md', '.json']:
                        title = title.replace(ext, '')
                    
                    book_data = {
                        "id": book_id,
                        "user_id": user_id,
                        "title": title,
                        "filename": filename,
                        "original_filename": filename,
                        "file_size": len(file_content),
                        "file_type": metadata.get('file_type', 'unknown'),
                        "format_details": metadata.get('format_details', 'Unknown Format'),
                        "page_count": metadata.get('page_count'),
                        "status": "processing"
                    }
                    
                    self.supabase_client.table("books").insert(book_data).execute()
                    logger.info(f"✅ Book metadata stored in Supabase: {book_id}")
                except Exception as e:
                    logger.error(f"❌ Error storing book in Supabase: {e}")
                    raise
            
            # Process with RAG if available
            if LANGCHAIN_AVAILABLE and self.text_splitter and self.embeddings and self.pinecone_index:
                try:
                    # Split text into chunks
                    documents = self.text_splitter.split_text(text)
                    logger.info(f"Split into {len(documents)} chunks")
                    
                    # Create embeddings and store in Pinecone
                    vectors_to_upsert = []
                    
                    for i, chunk in enumerate(documents):
                        # Generate embedding
                        embedding = self.embeddings.embed_query(chunk)
                        
                        # Create vector
                        vector_id = f"{book_id}_{i}"
                        vectors_to_upsert.append({
                            "id": vector_id,
                            "values": embedding,
                            "metadata": {
                                "book_id": book_id,
                                "user_id": user_id,
                                "filename": filename,
                                "chunk_index": i,
                                "text": chunk
                            }
                        })
                    
                    # Batch upsert to Pinecone
                    batch_size = 100
                    for i in range(0, len(vectors_to_upsert), batch_size):
                        batch = vectors_to_upsert[i:i + batch_size]
                        self.pinecone_index.upsert(vectors=batch)
                    
                    logger.info(f"✅ Stored {len(vectors_to_upsert)} vectors in Pinecone")
                    
                    # Update book status to processed
                    if self.supabase_client:
                        self.supabase_client.table("books").update({
                            "status": "processed",
                            "processed_date": datetime.now().isoformat()
                        }).eq("id", book_id).execute()
                        logger.info(f"✅ Book status updated to processed")
                        
                except Exception as e:
                    logger.error(f"Error in RAG processing: {e}")
                    # Update status to failed
                    if self.supabase_client:
                        self.supabase_client.table("books").update({
                            "status": "failed"
                        }).eq("id", book_id).execute()
                    raise
            
            logger.info(f"Successfully processed document: {filename} with ID: {book_id}")
            return book_id
            
        except Exception as e:
            logger.error(f"Error processing document: {e}")
            raise

    async def scrape_and_process_url(self, url: str, user_id: str, title: str = None) -> str:
        """Scrape content from URL and process it as a document"""
        if not WEB_SCRAPING_AVAILABLE:
            raise ValueError("Web scraping libraries not available")
        
        try:
            book_id = str(uuid.uuid4())
            
            # Validate URL
            if not validators.url(url):
                raise ValueError(f"Invalid URL: {url}")
            
            logger.info(f"Scraping URL: {url} for user: {user_id}")
            
            # Try multiple scraping methods
            text, metadata = await self._scrape_url_content(url)
            
            if not text or len(text.strip()) < 100:
                raise ValueError("Insufficient content extracted from URL")
            
            # Use provided title or extract from metadata
            document_title = title or metadata.get('title', urlparse(url).netloc)
            
            # Store in Supabase
            if self.supabase_client:
                try:
                    book_data = {
                        "id": book_id,
                        "user_id": user_id,
                        "title": document_title,
                        "filename": f"{document_title}.html",
                        "original_filename": url,
                        "file_size": len(text.encode()),
                        "file_type": "web",
                        "format_details": metadata.get('format_details', 'Web Content'),
                        "url_source": url,
                        "author": metadata.get('author'),
                        "publish_date": metadata.get('publish_date'),
                        "status": "processing"
                    }
                    
                    self.supabase_client.table("books").insert(book_data).execute()
                    logger.info(f"✅ Web content metadata stored in Supabase: {book_id}")
                except Exception as e:
                    logger.error(f"❌ Error storing web content in Supabase: {e}")
                    raise
            
            # Process with RAG if available
            if LANGCHAIN_AVAILABLE and self.text_splitter and self.embeddings and self.pinecone_index:
                try:
                    # Split text into chunks
                    documents = self.text_splitter.split_text(text)
                    logger.info(f"Split into {len(documents)} chunks")
                    
                    # Create embeddings and store in Pinecone
                    vectors_to_upsert = []
                    
                    for i, chunk in enumerate(documents):
                        # Generate embedding
                        embedding = self.embeddings.embed_query(chunk)
                        
                        # Create vector
                        vector_id = f"{book_id}_{i}"
                        vectors_to_upsert.append({
                            "id": vector_id,
                            "values": embedding,
                            "metadata": {
                                "book_id": book_id,
                                "user_id": user_id,
                                "filename": document_title,
                                "chunk_index": i,
                                "text": chunk,
                                "source_url": url
                            }
                        })
                    
                    # Batch upsert to Pinecone
                    batch_size = 100
                    for i in range(0, len(vectors_to_upsert), batch_size):
                        batch = vectors_to_upsert[i:i + batch_size]
                        self.pinecone_index.upsert(vectors=batch)
                    
                    logger.info(f"✅ Stored {len(vectors_to_upsert)} vectors in Pinecone")
                    
                    # Update status to processed
                    if self.supabase_client:
                        self.supabase_client.table("books").update({
                            "status": "processed",
                            "processed_date": datetime.now().isoformat()
                        }).eq("id", book_id).execute()
                        logger.info(f"✅ Web content status updated to processed")
                        
                except Exception as e:
                    logger.error(f"Error in RAG processing: {e}")
                    # Update status to failed
                    if self.supabase_client:
                        self.supabase_client.table("books").update({
                            "status": "failed"
                        }).eq("id", book_id).execute()
                    raise
            
            logger.info(f"Successfully processed web content from: {url} with ID: {book_id}")
            return book_id
            
        except Exception as e:
            logger.error(f"Error scraping and processing URL: {e}")
            raise

    async def _scrape_url_content(self, url: str) -> Tuple[str, Dict[str, Any]]:
        """Scrape content from URL using Firecrawl"""
        try:
            # Initialize Firecrawl
            firecrawl_api_key = os.getenv("FIRECRAWL_API_KEY")
            if not firecrawl_api_key:
                raise ValueError("FIRECRAWL_API_KEY not found in environment variables")
            
            app = FirecrawlApp(api_key=firecrawl_api_key)
            
            # Scrape with Firecrawl - get both markdown and HTML
            scrape_result = app.scrape_url(url, formats=['markdown', 'html'])
            
            if not scrape_result:
                raise ValueError(f"Firecrawl failed to scrape URL: {url}")
            
            # Extract content from the response object
            markdown_content = getattr(scrape_result, 'markdown', '') or ''
            html_content = getattr(scrape_result, 'html', '') or ''
            metadata = getattr(scrape_result, 'metadata', {}) or {}
            
            if not markdown_content or len(markdown_content.strip()) < 100:
                raise ValueError("Insufficient content extracted from URL")
            
            # Enhanced metadata with Firecrawl data
            enhanced_metadata = {
                'url': url,
                'title': metadata.get('title') if isinstance(metadata, dict) else None,
                'description': metadata.get('description') if isinstance(metadata, dict) else None,
                'language': metadata.get('language') if isinstance(metadata, dict) else None,
                'keywords': metadata.get('keywords') if isinstance(metadata, dict) else None,
                'author': metadata.get('ogTitle') if isinstance(metadata, dict) else None,  # Use ogTitle as author fallback
                'publish_date': None,  # Firecrawl doesn't provide publish date directly
                'extraction_method': 'firecrawl',
                'content_length': len(markdown_content),
                'format_details': 'Web Content (Firecrawl)',
                'status_code': metadata.get('statusCode') if isinstance(metadata, dict) else None,
                'source_url': metadata.get('sourceURL', url) if isinstance(metadata, dict) else url,
                'og_title': metadata.get('ogTitle') if isinstance(metadata, dict) else None,
                'og_description': metadata.get('ogDescription') if isinstance(metadata, dict) else None,
                'og_image': metadata.get('ogImage') if isinstance(metadata, dict) else None,
                'og_url': metadata.get('ogUrl') if isinstance(metadata, dict) else None,
                'robots': metadata.get('robots') if isinstance(metadata, dict) else None,
                'html_content': html_content  # Store HTML for potential future use
            }
            
            return markdown_content.strip(), enhanced_metadata
            
        except Exception as e:
            logger.error(f"Firecrawl scraping failed for {url}: {e}")
            raise ValueError(f"Firecrawl scraping failed: {e}")





    async def get_user_books(self, user_id: str) -> List[Book]:
        """Get all books for a specific user from Supabase"""
        try:
            if self.supabase_client:
                response = self.supabase_client.table("books").select("*").eq("user_id", user_id).order("upload_date", desc=True).execute()
                
                user_books = []
                for book_data in response.data:
                    book = Book(
                        id=book_data["id"],
                        title=book_data["title"],
                        filename=book_data["filename"],
                        original_filename=book_data["original_filename"],
                        user_id=book_data["user_id"],
                        upload_date=datetime.fromisoformat(book_data["upload_date"].replace('Z', '+00:00')),
                        file_size=book_data.get("file_size"),
                        page_count=book_data.get("page_count"),
                        status=book_data["status"],
                        processed_date=datetime.fromisoformat(book_data["processed_date"].replace('Z', '+00:00')) if book_data.get("processed_date") else None
                    )
                    user_books.append(book)
                
                return user_books
            else:
                logger.warning("Supabase not available")
                return []
        except Exception as e:
            logger.error(f"Error getting user books: {e}")
            return []

    async def chat_with_book(
        self, 
        book_id: str, 
        message: str, 
        user_id: str, 
        chat_history: List[ChatMessage] = None
    ) -> ChatResponse:
        """Chat with book using RAG"""
        try:
            # Get book from Supabase
            if not self.supabase_client:
                raise ValueError("Database not available")
                
            response = self.supabase_client.table("books").select("*").eq("id", book_id).eq("user_id", user_id).execute()
            
            if not response.data:
                raise ValueError("Book not found")
            
            book_data = response.data[0]
            book_title = book_data["title"]
            
            # Try RAG with Pinecone if available
            if (LANGCHAIN_AVAILABLE and self.pinecone_index and 
                self.embeddings and self.llm):
                try:
                    # Generate embedding for the question
                    question_embedding = self.embeddings.embed_query(message)
                    
                    # Search for relevant chunks in Pinecone
                    search_results = self.pinecone_index.query(
                        vector=question_embedding,
                        filter={"book_id": book_id},
                        top_k=4,
                        include_metadata=True
                    )
                    
                    if search_results.matches:
                        # Extract relevant text chunks
                        relevant_chunks = []
                        for match in search_results.matches:
                            chunk_text = match.metadata.get('text', '')
                            if chunk_text:
                                relevant_chunks.append(chunk_text)
                        
                        # Build context from relevant chunks
                        context = "\n\n".join(relevant_chunks[:3])
                        
                        # Build chat history context
                        history_context = ""
                        if chat_history:
                            recent_history = chat_history[-6:]  # Last 3 exchanges
                            for msg in recent_history:
                                history_context += f"{msg.role}: {msg.content}\n"
                        
                        # Create prompt for Gemini
                        prompt = f"""You are an AI tutor helping a student understand the book "{book_title}". 
Based on the following context from the book, answer the student's question clearly and helpfully.

Book Context:
{context}

{"Chat History:" + history_context if history_context else ""}

Student Question: {message}

Instructions:
- Answer based on the provided context from the book
- Be educational and clear
- If the context doesn't contain relevant information, say so
- Cite specific parts of the text when relevant
- Keep responses focused and helpful

Answer:"""
                        
                        # Get response from Gemini
                        ai_response = self.llm.invoke(prompt).content
                        
                        logger.info(f"✅ RAG response generated using {len(relevant_chunks)} chunks")
                        
                        # Save chat history
                        await self._save_chat_message(book_id, user_id, message, ai_response)
                        
                        return ChatResponse(
                            response=ai_response,
                            sources=[],
                            book_id=book_id
                        )
                        
                    else:
                        ai_response = f"I couldn't find specific information about '{message}' in the book '{book_title}'. Could you try rephrasing your question?"
                        
                        # Save chat history
                        await self._save_chat_message(book_id, user_id, message, ai_response)
                        
                        return ChatResponse(
                            response=ai_response,
                            sources=[],
                            book_id=book_id
                        )
                        
                except Exception as e:
                    logger.error(f"Error in RAG processing: {e}")
                    ai_response = f"I encountered an issue searching the book: {str(e)}"
                    await self._save_chat_message(book_id, user_id, message, ai_response)
                    
                    return ChatResponse(
                        response=ai_response,
                        sources=[],
                        book_id=book_id
                    )
            
            # Fallback response
            ai_response = f"I'm having trouble accessing the content of '{book_title}'. Please try uploading the book again."
            await self._save_chat_message(book_id, user_id, message, ai_response)
            
            return ChatResponse(
                response=ai_response,
                sources=[],
                book_id=book_id
            )
            
        except Exception as e:
            logger.error(f"Error in chat: {e}")
            raise

    async def _save_chat_message(self, book_id: str, user_id: str, user_message: str, ai_response: str):
        """Save chat messages to Supabase"""
        if not self.supabase_client:
            logger.warning("Supabase not available, chat history not saved")
            return
            
        try:
            # Save both user message and AI response
            messages_to_save = [
                {
                    "book_id": book_id,
                    "user_id": user_id,
                    "role": "user",
                    "content": user_message
                },
                {
                    "book_id": book_id,
                    "user_id": user_id,
                    "role": "assistant",
                    "content": ai_response
                }
            ]
            
            self.supabase_client.table("chat_histories").insert(messages_to_save).execute()
            logger.info(f"✅ Chat messages saved to Supabase for book: {book_id}")
        except Exception as e:
            logger.error(f"❌ Error saving chat messages to Supabase: {e}")

    async def get_chat_history(self, book_id: str, user_id: str) -> List[ChatMessage]:
        """Get chat history for a specific book from Supabase"""
        if not self.supabase_client:
            logger.warning("Supabase not available, returning empty chat history")
            return []
            
        try:
            response = self.supabase_client.table("chat_histories").select("*").eq("book_id", book_id).eq("user_id", user_id).order("created_at", desc=False).execute()
            
            chat_history = []
            for message_data in response.data:
                message = ChatMessage(
                    role=message_data["role"],
                    content=message_data["content"],
                    timestamp=datetime.fromisoformat(message_data["created_at"].replace('Z', '+00:00'))
                )
                chat_history.append(message)
            
            return chat_history
        except Exception as e:
            logger.error(f"❌ Error getting chat history from Supabase: {e}")
            return []

    async def delete_book(self, book_id: str, user_id: str):
        """Delete a book and its associated data"""
        try:
            # Delete from Supabase
            if self.supabase_client:
                self.supabase_client.table("books").delete().eq("id", book_id).eq("user_id", user_id).execute()
                logger.info(f"✅ Book deleted from Supabase: {book_id}")
            
            # Delete from Pinecone
            if self.pinecone_index:
                self.pinecone_index.delete(filter={"book_id": book_id})
                logger.info(f"✅ Book vectors deleted from Pinecone: {book_id}")
            
        except Exception as e:
            logger.error(f"Error deleting book: {e}")
            raise

    async def generate_mcqs(self, book_id: str, num_questions: int, difficulty: str, user_id: str) -> List[dict]:
        """Generate MCQs from a book using RAG and LLM"""
        try:
            if not self.initialized or not self.pinecone_index or not self.llm:
                raise Exception("RAG service not fully initialized")

            # Get relevant chunks from the book
            results = self.pinecone_index.query(
                vector=[0] * 768,  # Dummy vector to get all chunks
                filter={"book_id": book_id},
                top_k=10,
                include_metadata=True
            )

            # Combine chunks into context
            context = "\n".join([result.metadata.get("text", "") for result in results.matches])

            if not context:
                logger.warning(f"No content found for book_id: {book_id}")
                return self._generate_fallback_mcqs(num_questions, difficulty)

            # Generate MCQs using Gemini
            prompt = f"""Based on the following text, generate {num_questions} multiple choice questions at {difficulty} difficulty level.
            Each question should have 4 options (A, B, C, D) with exactly one correct answer.
            
Text: {context[:3000]}  
            
IMPORTANT: Return ONLY valid JSON in this exact format without any additional text or markdown:
            {{
                "mcqs": [
                    {{
                        "id": 1,
            "question": "Question text?",
                        "options": ["Option A", "Option B", "Option C", "Option D"],
            "correct_answer": 0,
                        "explanation": "Explanation of the correct answer"
                    }}
                ]
}}"""

            response = await self.llm.ainvoke(prompt)
            response_text = response.content.strip()
            
            # Clean the response text
            if response_text.startswith("```json"):
                response_text = response_text[7:]
            if response_text.endswith("```"):
                response_text = response_text[:-3]
            response_text = response_text.strip()

            # Parse JSON response
            try:
                mcq_data = json.loads(response_text)
                mcqs = mcq_data.get("mcqs", [])
                
                # Validate MCQ structure
                valid_mcqs = []
                for i, mcq in enumerate(mcqs):
                    if self._validate_mcq(mcq):
                        mcq["id"] = i + 1  # Ensure sequential IDs
                        valid_mcqs.append(mcq)
                
                if len(valid_mcqs) > 0:
                    logger.info(f"✅ Generated {len(valid_mcqs)} valid MCQs")
                    return valid_mcqs
                else:
                    logger.warning("No valid MCQs found in LLM response")
                    return self._generate_fallback_mcqs(num_questions, difficulty)
                    
            except json.JSONDecodeError as e:
                logger.error(f"Failed to parse LLM response as JSON: {e}")
                logger.error(f"Raw response: {response_text[:500]}...")
                return self._generate_fallback_mcqs(num_questions, difficulty)

        except Exception as e:
            logger.error(f"Error generating MCQs: {e}")
            return self._generate_fallback_mcqs(num_questions, difficulty)
    
    def _validate_mcq(self, mcq: dict) -> bool:
        """Validate MCQ structure"""
        required_fields = ["question", "options", "correct_answer", "explanation"]
        
        for field in required_fields:
            if field not in mcq:
                return False
        
        if not isinstance(mcq["options"], list) or len(mcq["options"]) != 4:
            return False
            
        if not isinstance(mcq["correct_answer"], int) or mcq["correct_answer"] < 0 or mcq["correct_answer"] > 3:
            return False
            
        return True
    
    def _generate_fallback_mcqs(self, num_questions: int, difficulty: str) -> List[dict]:
        """Generate fallback MCQs when LLM fails"""
        difficulties = {
            'easy': 'basic',
            'medium': 'intermediate', 
            'hard': 'advanced'
        }
        
        mcqs = []
        for i in range(num_questions):
            mcq = {
                "id": i + 1,
                "question": f"Sample {difficulties.get(difficulty, 'general')} question {i + 1} about the book content?",
                "options": [
                    f"Option A for question {i + 1}",
                    f"Option B for question {i + 1}",
                    f"Option C for question {i + 1}",
                    f"Option D for question {i + 1}"
                ],
                "correct_answer": 0,
                "explanation": f"This is a sample explanation for question {i + 1}. The actual MCQs would be generated from your book content."
            }
            mcqs.append(mcq)
        
        return mcqs 