import os
import io
import logging
import mimetypes
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import asyncio
import aiofiles

# Document processing imports
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
import xlrd
from pptx import Presentation
import markdown
import json
import csv

# Image processing
from PIL import Image
import pytesseract

# Audio processing
import speech_recognition as sr
from pydub import AudioSegment

# Text processing
import chardet
import magic

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Enhanced document processor supporting multiple file formats"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'word': ['.docx', '.doc'],
        'excel': ['.xlsx', '.xls', '.csv'],
        'powerpoint': ['.pptx', '.ppt'],
        'text': ['.txt', '.md', '.markdown', '.json'],
        'image': ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'],
        'audio': ['.wav', '.mp3', '.m4a', '.flac', '.aac'],
        'web': ['.html', '.htm']
    }
    
    def __init__(self):
        self.ocr_available = self._check_ocr_availability()
        self.speech_recognition_available = self._check_speech_recognition()
    
    def _check_ocr_availability(self) -> bool:
        """Check if OCR is available"""
        try:
            pytesseract.get_tesseract_version()
            return True
        except Exception:
            logger.warning("Tesseract OCR not available - image text extraction disabled")
            return False
    
    def _check_speech_recognition(self) -> bool:
        """Check if speech recognition is available"""
        try:
            sr.Recognizer()
            return True
        except Exception:
            logger.warning("Speech recognition not available - audio transcription disabled")
            return False
    
    def get_file_type(self, filename: str, content: bytes = None) -> str:
        """Determine file type from filename and content"""
        extension = Path(filename).suffix.lower()
        
        for format_type, extensions in self.SUPPORTED_FORMATS.items():
            if extension in extensions:
                return format_type
        
        # Try to detect MIME type if extension is unknown
        if content:
            try:
                mime_type = magic.from_buffer(content, mime=True)
                if mime_type.startswith('text/'):
                    return 'text'
                elif mime_type.startswith('image/'):
                    return 'image'
                elif mime_type.startswith('audio/'):
                    return 'audio'
                elif 'pdf' in mime_type:
                    return 'pdf'
            except Exception:
                pass
        
        return 'text'  # Default fallback
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported"""
        extension = Path(filename).suffix.lower()
        for extensions in self.SUPPORTED_FORMATS.values():
            if extension in extensions:
                return True
        return False
    
    async def extract_text(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from various file formats"""
        file_type = self.get_file_type(filename, content)
        metadata = {
            'filename': filename,
            'file_type': file_type,
            'file_size': len(content)
        }
        
        try:
            if file_type == 'pdf':
                text, pdf_metadata = self._extract_pdf_text(content)
                metadata.update(pdf_metadata)
            elif file_type == 'word':
                text, word_metadata = self._extract_word_text(content)
                metadata.update(word_metadata)
            elif file_type == 'excel':
                text, excel_metadata = self._extract_excel_text(content, filename)
                metadata.update(excel_metadata)
            elif file_type == 'powerpoint':
                text, ppt_metadata = self._extract_powerpoint_text(content)
                metadata.update(ppt_metadata)
            elif file_type == 'text':
                text, text_metadata = self._extract_text_content(content, filename)
                metadata.update(text_metadata)
            elif file_type == 'image':
                text, img_metadata = self._extract_image_text(content)
                metadata.update(img_metadata)
            elif file_type == 'audio':
                text, audio_metadata = await self._extract_audio_text(content, filename)
                metadata.update(audio_metadata)
            elif file_type == 'web':
                text, web_metadata = self._extract_html_text(content)
                metadata.update(web_metadata)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {e}")
            raise
    
    def _extract_pdf_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF files"""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            metadata = {
                'page_count': len(pdf_reader.pages),
                'format_details': 'PDF Document'
            }
            
            text = ""
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num + 1} ---\n"
                        text += page_text
                except Exception as e:
                    logger.warning(f"Could not extract text from page {page_num + 1}: {e}")
            
            if not text.strip():
                raise ValueError("No text could be extracted from PDF")
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing PDF: {e}")
    
    def _extract_word_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Word documents"""
        try:
            doc_file = io.BytesIO(content)
            doc = DocxDocument(doc_file)
            
            text = ""
            paragraph_count = 0
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
                    paragraph_count += 1
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text += cell.text + "\n"
            
            metadata = {
                'paragraph_count': paragraph_count,
                'table_count': len(doc.tables),
                'format_details': 'Word Document'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing Word document: {e}")
    
    def _extract_excel_text(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Excel files"""
        try:
            extension = Path(filename).suffix.lower()
            
            if extension == '.csv':
                return self._extract_csv_text(content)
            
            text = ""
            sheet_count = 0
            
            if extension == '.xlsx':
                workbook = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
                sheet_names = workbook.sheetnames
                
                for sheet_name in sheet_names:
                    sheet = workbook[sheet_name]
                    text += f"\n--- Sheet: {sheet_name} ---\n"
                    
                    for row in sheet.iter_rows():
                        row_text = []
                        for cell in row:
                            if cell.value is not None:
                                row_text.append(str(cell.value))
                        if row_text:
                            text += " | ".join(row_text) + "\n"
                    
                    sheet_count += 1
            
            elif extension == '.xls':
                workbook = xlrd.open_workbook(file_contents=content)
                
                for sheet_idx in range(workbook.nsheets):
                    sheet = workbook.sheet_by_index(sheet_idx)
                    text += f"\n--- Sheet: {sheet.name} ---\n"
                    
                    for row_idx in range(sheet.nrows):
                        row_text = []
                        for col_idx in range(sheet.ncols):
                            cell_value = sheet.cell_value(row_idx, col_idx)
                            if cell_value:
                                row_text.append(str(cell_value))
                        if row_text:
                            text += " | ".join(row_text) + "\n"
                    
                    sheet_count += 1
            
            metadata = {
                'sheet_count': sheet_count,
                'format_details': f'Excel Spreadsheet ({extension})'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing Excel file: {e}")
    
    def _extract_csv_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from CSV files"""
        try:
            # Detect encoding
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            
            text_content = content.decode(encoding)
            csv_reader = csv.reader(io.StringIO(text_content))
            
            text = ""
            row_count = 0
            
            for row in csv_reader:
                if row:  # Skip empty rows
                    text += " | ".join(str(cell) for cell in row) + "\n"
                    row_count += 1
            
            metadata = {
                'row_count': row_count,
                'encoding': encoding,
                'format_details': 'CSV File'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing CSV file: {e}")
    
    def _extract_powerpoint_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PowerPoint presentations"""
        try:
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            
            text = ""
            slide_count = 0
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                
                slide_count += 1
            
            metadata = {
                'slide_count': slide_count,
                'format_details': 'PowerPoint Presentation'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing PowerPoint file: {e}")
    
    def _extract_text_content(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from text-based files"""
        try:
            # Detect encoding
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            
            text = content.decode(encoding)
            extension = Path(filename).suffix.lower()
            
            metadata = {
                'encoding': encoding,
                'character_count': len(text),
                'line_count': len(text.splitlines()),
                'format_details': f'Text File ({extension})'
            }
            
            # Process markdown files
            if extension in ['.md', '.markdown']:
                # Convert markdown to plain text (optional: keep markdown formatting)
                metadata['format_details'] = 'Markdown Document'
            
            # Process JSON files
            elif extension == '.json':
                try:
                    json_data = json.loads(text)
                    # Convert JSON to readable text
                    text = json.dumps(json_data, indent=2)
                    metadata['format_details'] = 'JSON Document'
                except json.JSONDecodeError:
                    pass  # Keep as plain text
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing text file: {e}")
    
    def _extract_image_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from images using OCR"""
        if not self.ocr_available:
            raise ValueError("OCR not available - cannot extract text from images")
        
        try:
            image = Image.open(io.BytesIO(content))
            
            metadata = {
                'image_size': image.size,
                'image_mode': image.mode,
                'format_details': f'Image with OCR ({image.format})'
            }
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            if not text.strip():
                text = "[No text found in image]"
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing image: {e}")
    
    async def _extract_audio_text(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from audio files using speech recognition"""
        if not self.speech_recognition_available:
            raise ValueError("Speech recognition not available")
        
        try:
            # Save audio to temporary file
            temp_path = f"temp_audio_{hash(content)}.wav"
            
            # Convert to WAV format using pydub
            audio = AudioSegment.from_file(io.BytesIO(content))
            audio.export(temp_path, format="wav")
            
            # Perform speech recognition
            r = sr.Recognizer()
            with sr.AudioFile(temp_path) as source:
                audio_data = r.record(source)
                text = r.recognize_google(audio_data)
            
            # Cleanup
            if os.path.exists(temp_path):
                os.remove(temp_path)
            
            metadata = {
                'duration_seconds': len(audio) / 1000,
                'format_details': 'Audio Transcription'
            }
            
            return text, metadata
            
        except Exception as e:
            # Cleanup on error
            temp_path = f"temp_audio_{hash(content)}.wav"
            if os.path.exists(temp_path):
                os.remove(temp_path)
            raise ValueError(f"Error processing audio file: {e}")
    
    def _extract_html_text(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Extract text from HTML files"""
        try:
            from bs4 import BeautifulSoup
            import html2text
            
            # Detect encoding
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            
            html_content = content.decode(encoding)
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Convert to markdown-like text
            h = html2text.HTML2Text()
            h.ignore_links = False
            text = h.handle(html_content)
            
            metadata = {
                'encoding': encoding,
                'title': soup.title.string if soup.title else 'No title',
                'format_details': 'HTML Document'
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            raise ValueError(f"Error processing HTML file: {e}") 